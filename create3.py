import streamlit as st
from openai import OpenAI
import fitz  # PyMuPDF
from pptx import Presentation

st.set_page_config(page_title="AI 자기소개서 생성기", page_icon="🧑‍💼")
st.title("🧑‍💼 AI 자기소개서 에세이 생성기")

# 🔐 OpenAI API Key 입력
api_key = st.text_input("🔑 OpenAI API 키를 입력하세요", type="password")
if not api_key:
    st.warning("⚠️ OpenAI API 키를 입력해주세요.")
    st.stop()
client = OpenAI(api_key=api_key)

# 📌 기업 인재상 데이터
company_values = {
    "삼성전자": ["도전정신", "창의성", "글로벌 역량"],
    "LG전자": ["고객지향", "자율과 책임", "지속적 혁신"],
    "SK하이닉스": ["패기", "협력", "지속가능한 성장"],
    "현대자동차": ["도전", "소통과 협업", "고객 최우선"],
    "기아": ["혁신", "열정", "소통"],
    "카카오": ["유연한 사고", "문제 해결력", "소통 능력"],
    "네이버": ["기술 중심", "자율성", "글로벌 확장"],
    "롯데": ["도전", "창의", "책임감"],
    "포스코": ["고객 중심", "창의와 혁신", "신뢰와 협력"],
    "CJ": ["ONLYONE 정신", "창의성", "열정"],
    "한화": ["신뢰", "도전", "헌신"],
    "대한항공": ["책임감", "정직", "글로벌 역량"],
    "KT": ["도전", "혁신", "고객가치 중심"],
    "LG화학": ["책임감", "협업", "전문성"],
    "NHN": ["협업", "기술 성장", "유연성"],
    "쿠팡": ["고객 집착", "속도", "소유 의식"],
    "토스": ["문제 해결 중심", "책임", "신속한 실행"],
    "배달의민족(우아한형제들)": ["자율", "유쾌함", "사용자 중심"],
    "삼성바이오로직스": ["정직", "혁신", "책임의식"]
}

# 기업 선택
company = st.selectbox("📌 지원할 기업을 선택하세요", list(company_values.keys()))
if company:
    st.markdown(f"### 🏢 {company}의 인재상")
    for v in company_values[company]:
        st.markdown(f"- {v}")

# 📎 파일 업로드 (PDF/PPTX)
st.header("📎 PDF 또는 PPTX 참고자료 업로드 (선택)")
uploaded_file = st.file_uploader("자소서 참고용 파일을 업로드하세요", type=["pdf", "pptx"])

# 파일 내용 추출
def extract_text_from_pdf(file):
    text = ""
    with fitz.open(stream=file.read(), filetype="pdf") as doc:
        for page in doc:
            text += page.get_text()
    return text

def extract_text_from_pptx(file):
    prs = Presentation(file)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

pdf_text = ""
if uploaded_file:
    if uploaded_file.type == "application/pdf":
        pdf_text = extract_text_from_pdf(uploaded_file)
        st.success("✅ PDF에서 텍스트 추출 완료.")
    elif uploaded_file.type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
        pdf_text = extract_text_from_pptx(uploaded_file)
        st.success("✅ PPTX에서 텍스트 추출 완료.")
    else:
        st.error("❌ 지원하지 않는 파일 형식입니다.")

# 자기소개서 입력
st.header("📝 자기소개서 작성 입력")
reason = st.text_area("1. 지원 동기", height=100)
background = st.text_area("2. 성장 과정", height=100)
experience = st.text_area("3. 직무 관련 경험", height=100)

# 자기소개서 생성 함수
def generate_cover_letter(reason, background, experience, company, pdf_text=""):
    traits = ", ".join(company_values[company])
    prompt = f"""
아래 내용을 바탕으로 자연스럽고 진솔한 에세이 형식의 자기소개서를 작성해 주세요.
각 항목은 하나의 문단으로 만들고, 문단과 문단 사이에는 부드러운 연결 문장을 넣어 글의 흐름이 자연스럽게 이어지도록 해주세요.
너무 딱딱하거나 공식적인 표현보다는 개인적인 경험과 감정을 담아 진짜 이야기를 듣는 느낌이 들게 해주세요.

[지원 기업]: {company}
[기업 인재상]: {traits}

[지원 동기]
{reason}

[성장 과정]
{background}

[직무 경험]
{experience}
"""
    if pdf_text.strip():
        prompt += f"\n[PDF/PPTX 참고 내용]\n{pdf_text.strip()}\n\n이 내용도 참고하여 자기소개서를 작성해주세요.\n"

    response = client.chat.completions.create(
        model="gpt-4",
        messages=[{"role": "user", "content": prompt}],
        temperature=0.8
    )
    return response.choices[0].message.content.strip()

# 면접 질문 생성 함수
def generate_interview_questions(cover_letter_text, company):
    prompt = f"""
다음은 {company}에 지원한 자기소개서입니다. 이 자기소개서를 기반으로 실제 면접에서 나올 수 있는 심도 있는 질문 5개를 작성해주세요.
- 너무 일반적인 질문은 제외해주세요. (예: "자기소개 해보세요")
- 자기소개서에 언급된 경험, 가치관, 동기 등을 구체적으로 묻는 방식으로 만들어주세요.

[자기소개서]
{cover_letter_text}
"""
    response = client.chat.completions.create(
        model="gpt-4",
        messages=[{"role": "user", "content": prompt}],
        temperature=0.7
    )
    return response.choices[0].message.content.strip()

# 버튼 클릭 시 실행
if st.button("🚀 자기소개서 + 면접 질문 생성"):
    if not (reason and background and experience):
        st.error("❗ 모든 입력란을 채워주세요.")
    else:
        with st.spinner("자기소개서를 작성 중입니다..."):
            cover_letter = generate_cover_letter(reason, background, experience, company, pdf_text)
        st.subheader("📄 생성된 자기소개서")
        st.write(cover_letter)
        st.download_button("📥 자기소개서 다운로드", cover_letter, file_name="cover_letter.txt")

        with st.spinner("면접 질문을 생성 중입니다..."):
            interview_questions = generate_interview_questions(cover_letter, company)
        st.subheader("💬 예상 면접 질문")
        st.write(interview_questions)
